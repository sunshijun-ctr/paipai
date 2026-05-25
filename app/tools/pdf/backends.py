"""
PDF and text file extraction.
Features:
- Full text extraction (cleaned)
- Heuristic section detection (works well for CS/ML papers)
- Per-page/slide text for fine-grained retrieval
- Basic metadata extraction
依赖 (PDF): pymupdf
依赖 (PPTX): python-pptx
"""
import logging
import os
import re
from typing import Any, Optional

logger = logging.getLogger(__name__)


# ── Section header patterns ────────────────────────────────────────────────
# Maps canonical section name → list of lowercase patterns to match against

_SECTION_MAP: dict[str, list[str]] = {
    "abstract":      ["abstract"],
    "introduction":  ["introduction"],
    "related_work":  ["related work", "related works", "background", "prior work", "literature review"],
    "method":        ["method", "methods", "methodology", "approach", "proposed method",
                      "our approach", "proposed approach", "framework", "architecture"],
    "experiments":   ["experiment", "experiments", "experimental setup", "experimental results",
                      "evaluation", "results", "result", "benchmarks"],
    "discussion":    ["discussion", "analysis", "ablation study", "ablation"],
    "conclusion":    ["conclusion", "conclusions", "concluding remarks", "summary and conclusion"],
    "references":    ["references", "bibliography"],
}

# Flat lookup: pattern → section name
_PATTERN_TO_SECTION: dict[str, str] = {
    pat: sec for sec, pats in _SECTION_MAP.items() for pat in pats
}

_NUMBER_PREFIX = re.compile(r"^(?:\d+[\.\s]+|[ivxIVX]+[\.\s]+)")


def _normalize_header(line: str) -> str:
    """Strip leading numbering and lowercase for matching."""
    return _NUMBER_PREFIX.sub("", line.strip()).strip().lower()


def _is_section_header(line: str) -> Optional[str]:
    """
    Return the canonical section name if the line looks like a section header,
    otherwise None.
    A line qualifies when:
      - It's short (< 80 chars)
      - After stripping numbering it matches a known pattern
    """
    stripped = line.strip()
    if not stripped or len(stripped) > 80:
        return None
    normalized = _normalize_header(stripped)
    if normalized in _PATTERN_TO_SECTION:
        return _PATTERN_TO_SECTION[normalized]
    # Also try prefix match (e.g. "Experiments and Results" → "experiments")
    for pat, sec in _PATTERN_TO_SECTION.items():
        if normalized.startswith(pat):
            return sec
    return None


# ── Text cleaning ──────────────────────────────────────────────────────────

def _clean_text(text: str) -> str:
    # Fix hyphenated line breaks: "word-\nword" → "wordword"
    text = re.sub(r"-\n(\w)", r"\1", text)
    # Collapse multiple blank lines → one
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip trailing whitespace per line
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return text.strip()


def _paper_id(title_or_path: str) -> str:
    base = os.path.splitext(os.path.basename(title_or_path))[0] or title_or_path
    return re.sub(r"[^a-z0-9]+", "_", base.lower()).strip("_")[:64]


def _infer_chunk_type(text: str, section: str = "") -> str:
    sample = (section + "\n" + text[:400]).lower()
    if "|" in text and re.search(r"\n\s*\|?\s*[-:]{3,}", text):
        return "table"
    if re.search(r"\b(table|tab\.)\s*\d+", sample):
        return "table"
    if re.search(r"\b(fig\.|figure)\s*\d+", sample):
        return "figure"
    return "text"


def _split_markdown_sections(markdown: str) -> dict[str, str]:
    sections: dict[str, list[str]] = {}
    current = "body"
    sections[current] = []

    for line in markdown.splitlines():
        m = re.match(r"^\s{0,3}#{1,6}\s+(.+?)\s*$", line)
        if m:
            header = _normalize_header(m.group(1))
            current = _canonical_section(header) or _slug_section(header)
            count = 2
            base = current
            while current in sections:
                current = f"{base}_{count}"
                count += 1
            sections[current] = []
            continue
        sections[current].append(line)

    cleaned = {
        name: _clean_text("\n".join(lines))
        for name, lines in sections.items()
        if _clean_text("\n".join(lines))
    }
    return cleaned or _split_into_sections(markdown)


def _canonical_section(header: str) -> Optional[str]:
    if header in _PATTERN_TO_SECTION:
        return _PATTERN_TO_SECTION[header]
    for pat, sec in _PATTERN_TO_SECTION.items():
        if header.startswith(pat):
            return sec
    return None


def _slug_section(header: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "_", header.lower()).strip("_")
    return slug[:48] or "section"


def _build_rag_chunks(
    sections: dict[str, str],
    *,
    source: str,
    parser: str,
    pages_by_section: dict[str, int] | None = None,
) -> list[dict[str, Any]]:
    paper_id = _paper_id(source)
    chunks: list[dict[str, Any]] = []
    for section, text in sections.items():
        text = text.strip()
        if not text:
            continue
        chunks.append({
            "text": text,
            "metadata": {
                "paper_id": paper_id,
                "section": section,
                "page": (pages_by_section or {}).get(section, 0),
                "chunk_type": _infer_chunk_type(text, section),
                "parser": parser,
            },
        })
    return chunks


# ── Section splitting ──────────────────────────────────────────────────────

def _split_into_sections(full_text: str) -> dict[str, str]:
    """
    Heuristically split the text into named sections.
    Unknown sections are grouped under 'body' or numbered headers.
    """
    lines = full_text.splitlines()
    sections: dict[str, list[str]] = {}
    current_section = "preamble"
    sections[current_section] = []

    for line in lines:
        sec = _is_section_header(line)
        if sec:
            # Deduplicate section names (e.g. two "experiments" sections)
            key = sec
            count = 2
            while key in sections:
                key = f"{sec}_{count}"
                count += 1
            current_section = key
            sections[current_section] = []
        else:
            sections[current_section].append(line)

    # Join and clean each section; drop empty ones
    return {
        name: "\n".join(lines).strip()
        for name, lines in sections.items()
        if "\n".join(lines).strip()
    }


# ── Main extraction ────────────────────────────────────────────────────────

def extract(local_path: str) -> dict[str, Any]:
    """
    Extract text and structure from a PDF file.

    Returns:
        full_text: cleaned full document text
        sections:  dict of section_name → text
        pages:     list of per-page text (for chunking / RAG)
        page_count: int
        metadata:  dict from PDF metadata (title, author, etc.)
    """
    try:
        import fitz  # pymupdf
    except ImportError as e:
        raise ImportError("pymupdf is required: pip install pymupdf") from e

    try:
        doc = fitz.open(local_path)
    except Exception as e:
        raise ValueError(f"Cannot open PDF: {local_path}") from e

    pages_text: list[str] = []
    for page in doc:
        pages_text.append(page.get_text("text"))

    raw_full = "\n".join(pages_text)
    full_text = _clean_text(raw_full)
    sections = _split_into_sections(full_text)

    metadata: dict[str, Any] = {}
    try:
        meta = doc.metadata or {}
        metadata = {k: v for k, v in meta.items() if v}
    except Exception:
        pass

    doc.close()

    return {
        "full_text": full_text,
        "sections": sections,
        "pages": [_clean_text(p) for p in pages_text],
        "page_count": len(pages_text),
        "metadata": metadata,
        "rag_chunks": _build_rag_chunks(sections, source=local_path, parser="pymupdf"),
    }


def extract_llama_index(local_path: str) -> dict[str, Any]:
    """Extract a PDF through LlamaParse/LlamaIndex and return the project shape.

    The RAG-facing format remains compatible with the existing architecture:
    full_text / sections / pages / metadata are preserved, and rag_chunks adds
    section-aware text/table/figure chunks for indexing.
    """
    from app.config.settings import settings

    api_key = settings.llama_cloud_api_key or os.environ.get("LLAMA_CLOUD_API_KEY")
    if not api_key:
        raise ImportError("LLAMA_CLOUD_API_KEY is required for LlamaParse")

    try:
        from llama_index.readers.llama_parse import LlamaParse
    except ImportError:
        try:
            from llama_parse import LlamaParse
        except ImportError as exc:
            raise ImportError(
                "llama-index LlamaParse support is required: pip install llama-index llama-parse"
            ) from exc

    parser = LlamaParse(
        api_key=api_key,
        result_type=settings.llama_parse_result_type or "markdown",
    )
    documents = parser.load_data(local_path)

    pages: list[str] = []
    page_numbers: list[int] = []
    metadata: dict[str, Any] = {"parser": "llama_parse"}
    for idx, doc in enumerate(documents, start=1):
        text = getattr(doc, "text", "") or ""
        if not text.strip():
            continue
        meta = getattr(doc, "metadata", {}) or {}
        page = meta.get("page_label") or meta.get("page") or idx
        try:
            page_num = int(page)
        except Exception:
            page_num = idx
        page_numbers.append(page_num)
        pages.append(_clean_text(text))
        for key in ("title", "author", "source"):
            if meta.get(key) and key not in metadata:
                metadata[key] = meta[key]

    full_text = _clean_text("\n\n".join(pages))
    sections = _split_markdown_sections(full_text)

    pages_by_section: dict[str, int] = {}
    for section, section_text in sections.items():
        for page_num, page_text in zip(page_numbers, pages):
            if section_text[:120] and section_text[:120] in page_text:
                pages_by_section[section] = page_num
                break

    return {
        "full_text": full_text,
        "sections": sections,
        "pages": pages,
        "page_count": len(pages),
        "metadata": metadata,
        "rag_chunks": _build_rag_chunks(
            sections,
            source=local_path,
            parser="llama_parse",
            pages_by_section=pages_by_section,
        ),
    }


# ── Plain-text extraction ──────────────────────────────────────────────────

def _extract_text_file(local_path: str) -> dict[str, Any]:
    """Read a plain-text or Markdown file and return the same shape as extract()."""
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            with open(local_path, encoding=enc) as f:
                raw = f.read()
            break
        except UnicodeDecodeError:
            continue
    else:
        raise ValueError(f"Cannot decode text file: {local_path}")

    full_text = _clean_text(raw)
    return {
        "full_text": full_text,
        "sections": {},
        "pages": [full_text],
        "page_count": 1,
        "metadata": {},
        "rag_chunks": _build_rag_chunks({"body": full_text}, source=local_path, parser="text"),
    }


_TEXT_EXTS = {".txt", ".md", ".text", ".rst"}
_PPTX_EXTS = {".pptx"}


def _shape_text(shape) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = (shape.text or "").strip()
        if text:
            texts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            row_text = " | ".join((cell.text or "").strip() for cell in row.cells)
            if row_text.strip(" |"):
                texts.append(row_text)
    if getattr(shape, "shape_type", None) == 6:  # GROUP
        for child in shape.shapes:
            texts.extend(_shape_text(child))
    return texts


def _extract_pptx_file(local_path: str) -> dict[str, Any]:
    """Extract text from a PowerPoint .pptx file into slide sections."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ImportError("python-pptx is required: pip install python-pptx") from e

    prs = Presentation(local_path)
    slides: list[str] = []
    sections: dict[str, str] = {}

    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"Slide {idx}"]
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))

        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes:\n{notes}")
        except Exception:
            pass

        slide_text = _clean_text("\n".join(p for p in parts if p.strip()))
        slides.append(slide_text)
        if len(slide_text) >= 20:
            sections[f"slide_{idx:03d}"] = slide_text

    core = prs.core_properties
    metadata = {
        "title": core.title or "",
        "author": core.author or "",
        "subject": core.subject or "",
        "keywords": core.keywords or "",
        "created": core.created.isoformat() if core.created else "",
        "modified": core.modified.isoformat() if core.modified else "",
    }
    metadata = {k: v for k, v in metadata.items() if v}

    return {
        "full_text": _clean_text("\n\n".join(slides)),
        "sections": sections,
        "pages": slides,
        "page_count": len(slides),
        "metadata": metadata,
        "rag_chunks": _build_rag_chunks(sections, source=local_path, parser="pptx"),
    }


def extract_any(local_path: str) -> dict[str, Any]:
    """Dispatch to the right extractor based on file extension.

    PDFs go through LlamaParse first when configured. We fall back to
    PyMuPDF when LlamaParse either raises OR returns empty content
    (the latter happens on quota exhaustion, server-side parse errors,
    or oversized files — none of which raise on the client). Without
    the empty-result check, large PDFs silently land in the library
    with 0 chunks indexed.
    """
    ext = os.path.splitext(local_path)[1].lower()
    if ext in _TEXT_EXTS:
        return _extract_text_file(local_path)
    if ext in _PPTX_EXTS:
        return _extract_pptx_file(local_path)
    from app.config.settings import settings
    if settings.use_llama_parse:
        try:
            data = extract_llama_index(local_path)
            if _has_content(data):
                logger.info("PDF extracted via LlamaParse: %s", os.path.basename(local_path))
                return data
            logger.warning(
                "LlamaParse returned empty for '%s' — falling back to PyMuPDF "
                "(likely quota / size / upstream parse error)",
                os.path.basename(local_path),
            )
        except Exception as exc:
            logger.warning(
                "LlamaParse extraction failed for '%s' (%s) — falling back to PyMuPDF",
                os.path.basename(local_path), exc,
            )
    result = extract(local_path)
    logger.info(
        "PDF extracted via PyMuPDF: %s (%d chars, %d sections, %d chunks)",
        os.path.basename(local_path),
        len(result.get("full_text") or ""),
        len(result.get("sections") or {}),
        len(result.get("rag_chunks") or []),
    )
    return result


def _has_content(data: dict[str, Any]) -> bool:
    """A non-empty extraction has at least one of: full_text, any section,
    or any rag_chunk. Empty across all three means the extractor produced
    nothing usable."""
    if (data.get("full_text") or "").strip():
        return True
    sections = data.get("sections") or {}
    if any((v or "").strip() for v in sections.values()):
        return True
    return bool(data.get("rag_chunks"))
